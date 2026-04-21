# Step 1 — auto-install dependencies
import subprocess, sys
for pkg in ['python-pptx', 'Pillow']:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', pkg, '-q'])

# Step 2 — imports
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage

# Step 3 — constants
NAVY   = RGBColor(0x0B, 0x16, 0x26)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MINT   = RGBColor(0x3F, 0xFF, 0xAB)
GREY   = RGBColor(0xE5, 0xE7, 0xEB)
LGREY  = RGBColor(0x9C, 0xA3, 0xAF)

SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

TMP = 'tmp_images'
os.makedirs(TMP, exist_ok=True)

# Step 4 — helper: set background
def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

# Step 5 — helper: add label (small mint caps above title)
def add_label(slide, text):
    txb = slide.shapes.add_textbox(Cm(1.8), Cm(0.5), Cm(18), Cm(0.8))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = MINT
    run.font.name = 'Arial'

# Step 6 — helper: add title
def add_title(slide, text, has_image=True):
    max_w = Cm(33.87 * 0.55) if has_image else Cm(33.87 * 0.90)
    txb = slide.shapes.add_textbox(Cm(1.8), Cm(1.2), max_w, Cm(3.5))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    size = Pt(26) if len(text) > 45 else Pt(32)
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Arial'

# Step 7 — helper: add subtitle
def add_subtitle(slide, text, top_cm=3.0):
    txb = slide.shapes.add_textbox(Cm(1.8), Cm(top_cm), Cm(18), Cm(1.5))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(20)
    run.font.color.rgb = GREY
    run.font.name = 'Arial'

# Step 8 — helper: add body text (plain paragraph)
def add_body(slide, text, top_cm=4.0, width_cm=18):
    txb = slide.shapes.add_textbox(Cm(1.8), Cm(top_cm), Cm(width_cm), Cm(8))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(15)
    run.font.color.rgb = GREY
    run.font.name = 'Arial'

# Step 9 — helper: add bullets with mint marker
def add_bullets(slide, bullets, top_cm=4.2, width_cm=17.5):
    txb = slide.shapes.add_textbox(Cm(1.8), Cm(top_cm), Cm(width_cm), Cm(13))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        # mint marker run
        r1 = p.add_run()
        r1.text = '\u25b8 '
        r1.font.size = Pt(15)
        r1.font.color.rgb = MINT
        r1.font.name = 'Arial'
        # bullet text run
        r2 = p.add_run()
        r2.text = bullet
        r2.font.size = Pt(15)
        r2.font.color.rgb = GREY
        r2.font.name = 'Arial'

# Step 10 — helper: add quote
def add_quote(slide, text, top_cm=3.2, width_cm=17.5):
    txb = slide.shapes.add_textbox(Cm(1.8), Cm(top_cm), Cm(width_cm), Cm(3))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(17)
    run.font.italic = True
    run.font.color.rgb = MINT
    run.font.name = 'Arial'

# Step 11 — helper: convert webp to png, return usable path
def resolve_image(path):
    if not os.path.exists(path):
        return None
    if path.lower().endswith('.webp'):
        fname = os.path.splitext(os.path.basename(path))[0] + '.png'
        out = os.path.join(TMP, fname)
        if not os.path.exists(out):
            PILImage.open(path).convert('RGBA').save(out, 'PNG')
        return out
    return path

# Step 12 — helper: add image proportionally inside a zone
def add_image_proportional(slide, image_path, zone_left_cm, zone_top_cm, zone_w_cm, zone_h_cm):
    path = resolve_image(image_path)
    if path is None:
        print(f'  \u2717 Image not found: {image_path}')
        return
    with PILImage.open(path) as img:
        iw, ih = img.size
    zw = Cm(zone_w_cm)
    zh = Cm(zone_h_cm)
    scale = min(zw / iw, zh / ih)
    fw = iw * scale
    fh = ih * scale
    ox = (zw - fw) / 2
    oy = (zh - fh) / 2
    slide.shapes.add_picture(path, Cm(zone_left_cm) + ox, Cm(zone_top_cm) + oy, fw, fh)
    print(f'  \u2713 Image: {os.path.basename(image_path)}')

# Step 13 — standard right zone constants
RZ_LEFT  = 33.87 * 0.57   # cm
RZ_TOP   = 1.2             # cm
RZ_W     = 33.87 * 0.40   # cm
RZ_H     = 19.05 - 2.4    # cm

# Two-image zone constants
TZ1_LEFT = 33.87 * 0.57
TZ2_LEFT = 33.87 * 0.78
TZ_TOP   = 3.5
TZ_W     = 33.87 * 0.19
TZ_H     = 12.0

# ── CREATE PRESENTATION ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # completely blank layout

results = []

def new_slide(title_text):
    s = prs.slides.add_slide(blank)
    set_bg(s)
    results.append(title_text)
    return s

# ── SLIDE 1 — Cover ──────────────────────────────────────────────────────────
s = new_slide('Cover')
add_label(s, 'APRIL 2026 — TRAILVOLT')
add_title(s, 'TrailVolt', has_image=True)
add_subtitle(s, 'One Pack. Every Horizon.', top_cm=3.0)
add_body(s, 'How I used Claude Code + GitHub to build a premium e-commerce site for a concept-driven product \u2014 and why a boxed solution was never an option.', top_cm=4.2, width_cm=17)
add_image_proportional(s, 'assets/images/hero-pack-main.webp', 33.87*0.55, 0, 33.87*0.45, 19.05)

# ── SLIDE 2 — Why Not Shopify ────────────────────────────────────────────────
s = new_slide('Why Not Shopify or Wix')
add_label(s, 'THE TOOL QUESTION')
add_title(s, 'A premium product needs a premium digital environment', has_image=True)
add_bullets(s, [
    'Off-the-shelf platforms sell known categories well \u2014 they fail at explaining new product logic',
    'TrailVolt is a system, not a SKU \u2014 it cannot live inside a standard product card',
    'I needed a site that is itself part of the product, not just a storefront in front of it',
    'Claude Code + GitHub gave full control over architecture, UX flow, visual language and storytelling',
    "The result: a site that doesn't sell a backpack \u2014 it communicates a philosophy",
])
add_image_proportional(s, 'CCP/Front Page.png', RZ_LEFT, RZ_TOP, RZ_W, RZ_H)

# ── SLIDE 3 — The Problem ────────────────────────────────────────────────────
s = new_slide('The Problem')
add_label(s, 'THE ORIGIN')
add_title(s, 'Four bags. One person. One endless reorganisation.', has_image=True)
add_bullets(s, [
    'Laptop bag for work. Weekend duffel. Trail pack. Gym bag. Four bags, four systems',
    'Every Sunday: repack. Every Friday: repack again. Wrong shoes. Forgotten chargers',
    'Expensive backpacks failed the same way \u2014 designed for one type of journey',
    'Dublin Airport. Thursday morning. A man ahead unclips a module from his pack. The idea.',
    "What if the module wasn't just the bag shape \u2014 but the accessories too?",
])
add_image_proportional(s, 'CCP/3 tipe bags.png', RZ_LEFT, RZ_TOP, RZ_W, RZ_H)

# ── SLIDE 4 — The Idea ───────────────────────────────────────────────────────
s = new_slide('The Idea')
add_label(s, 'THE CONCEPT')
add_title(s, 'Not a backpack. A system.', has_image=True)
add_quote(s, '"1 base pack + 8 snap-on modules = everything you need"', top_cm=3.2)
add_bullets(s, [
    'One base pack in three configurations: Everyday 22L, Travel 30L, Trail 28L',
    'Eight cross-compatible snap-on modules \u2014 all work with all three lines',
    'Stainless steel attachment rail on the front face \u2014 visible proof of modularity',
    'Replaces four bags with one ecosystem',
    'From \u20ac159. 5-year warranty. Lifetime repair programme.',
], top_cm=6.0)
add_image_proportional(s, 'assets/images/modular-system-diagram.webp', RZ_LEFT, RZ_TOP, RZ_W, RZ_H)

# ── SLIDE 5 — Three Lines ────────────────────────────────────────────────────
s = new_slide('Three Lines')
add_label(s, 'THE ARCHITECTURE')
add_title(s, 'Three contexts. One philosophy.', has_image=False)

col_w = 33.87 / 3
col_data = [
    ('EVERYDAY 22L', 'City, office, gym, evening\n600D Ballistic Nylon / Core Black'),
    ('TRAVEL 30L',   'Airport, 4 nights, no checked bag\nIATA-compliant / Slate Grey'),
    ('TRAIL 28L',    'Wicklow, Wild Atlantic Way, Pyrenees\n500D Cordura / Forest Green'),
]
imgs_3 = ['assets/images/line-card-everyday.webp', 'assets/images/line-card-travel.webp', 'assets/images/line-card-trail.webp']
for i, (lbl, txt) in enumerate(col_data):
    x = i * col_w + 0.3
    # mint column label
    txb = s.shapes.add_textbox(Cm(x), Cm(3.2), Cm(col_w - 0.4), Cm(0.8))
    tf = txb.text_frame; p = tf.paragraphs[0]; r = p.add_run()
    r.text = lbl; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = MINT; r.font.name = 'Arial'
    # subtext
    txb2 = s.shapes.add_textbox(Cm(x), Cm(4.2), Cm(col_w - 0.4), Cm(2.5))
    tf2 = txb2.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]; r2 = p2.add_run()
    r2.text = txt; r2.font.size = Pt(13); r2.font.color.rgb = GREY; r2.font.name = 'Arial'
    # image
    add_image_proportional(s, imgs_3[i], x, 7.0, col_w - 0.4, 10.5)

# ── SLIDE 6 — Modular Ecosystem ──────────────────────────────────────────────
s = new_slide('Modular Ecosystem')
add_label(s, 'THE MODULES')
add_title(s, '8 modules. One rail. Zero compromise.', has_image=True)
add_bullets(s, [
    'Chest Phone Module \u2014 instant phone access, no bag opening required (\u20ac29)',
    'Tech Organizer \u2014 cables, SD cards, earbuds accessible on the outside (\u20ac35)',
    'Power Bank 10,000mAh \u2014 USB-A + USB-C charging while you move (\u20ac39)',
    'Shoe Module \u2014 ventilated isolated compartment for muddy gear (\u20ac32)',
    'Hidden Document Sleeve \u2014 RFID-blocking, passport + boarding pass (\u20ac19)',
    'Every module works with all three pack lines \u2014 no compatibility concerns',
])
add_image_proportional(s, 'assets/images/module-phone.webp', TZ1_LEFT, TZ_TOP, TZ_W, TZ_H)
add_image_proportional(s, 'assets/images/module-tech.webp',  TZ2_LEFT, TZ_TOP, TZ_W, TZ_H)

# ── SLIDE 7 — Homepage Strategy ──────────────────────────────────────────────
s = new_slide('Homepage Strategy')
add_label(s, 'UX ARCHITECTURE')
add_title(s, "The homepage is not a storefront. It's a product onboarding.", has_image=True)
add_bullets(s, [
    'Hero: pack large, dark background, mint accent \u2014 premium signal in under 2 seconds',
    'Trust strip immediately below hero: free shipping, warranty, returns, materials',
    'Press bar: WIRED, Gear Patrol, Monocle, Uncrate \u2014 authority without explanation',
    'Stats bar: 3 lines, 8 modules, 16" laptop, from \u20ac159 \u2014 the system in 4 numbers',
    'Line cards: three use contexts so the visitor recognises their own life',
    'Everything above the first scroll answers: is this a serious product?',
])
add_image_proportional(s, 'CCP/Front Page.png', RZ_LEFT, RZ_TOP, RZ_W, RZ_H)

# ── SLIDE 8 — Configurator ───────────────────────────────────────────────────
s = new_slide('Configurator')
add_label(s, 'KEY UX ELEMENT')
add_title(s, "The configurator is not an order form. It's the system interface.", has_image=True)
add_bullets(s, [
    '4 steps: line \u2192 colour \u2192 modules \u2192 build review',
    'Live price updates with every module selection',
    '"All modules are cross-compatible" \u2014 stated directly in step 3',
    'The user is not selecting a product \u2014 they are building their own kit',
    'This is the moment the visitor understands they are dealing with a system',
])
add_image_proportional(s, 'CCP/Configure core pack.png',        TZ1_LEFT, TZ_TOP, TZ_W, TZ_H)
add_image_proportional(s, 'CCP/Configure core back step 3.png', TZ2_LEFT, TZ_TOP, TZ_W, TZ_H)

# ── SLIDE 9 — Modules Page ───────────────────────────────────────────────────
s = new_slide('Modules Page')
add_label(s, 'PRODUCT DEPTH')
add_title(s, 'Modules are a full product \u2014 not an accessories section.', has_image=True)
add_bullets(s, [
    'Each module has its own article: photo, description, spec sheet, weight',
    'Page logic: ecosystem first, then individual module detail',
    'Feature strip at top: snap in 3 seconds / 100% interchangeable / no tools required',
    'Quick Add to Cart directly from the page \u2014 no configurator required',
    'Same logic Apple uses for AirPods: part of an ecosystem, not just an add-on',
])
add_image_proportional(s, 'assets/images/module-shoe.webp',  TZ1_LEFT, TZ_TOP, TZ_W, TZ_H)
add_image_proportional(s, 'assets/images/module-power.webp', TZ2_LEFT, TZ_TOP, TZ_W, TZ_H)

# ── SLIDE 10 — Founder Story ─────────────────────────────────────────────────
s = new_slide('Founder Story')
add_label(s, 'TRUST ARCHITECTURE')
add_title(s, 'Built Out of Frustration. Perfected by Obsession.', has_image=True)
add_quote(s, '"I wanted one pack that could do everything \u2014 not because I was lazy, but because the right system should be that capable."', top_cm=3.0)
# attribution
txb = s.shapes.add_textbox(Cm(1.8), Cm(5.4), Cm(17), Cm(0.7))
tf = txb.text_frame; p = tf.paragraphs[0]; r = p.add_run()
r.text = '\u2014 Ruslan, Founder'; r.font.size = Pt(13); r.font.color.rgb = LGREY; r.font.name = 'Arial'; r.font.italic = True
add_bullets(s, [
    'The founder story is not a marketing device \u2014 it explains why the product exists',
    'Real problem, real solution, real product tested in real conditions',
    '6+ rounds of physical prototyping. Dublin, Wicklow Mountains, airports',
    'A visible founder creates trust that no brand voice can manufacture',
], top_cm=6.5)
add_image_proportional(s, 'CCP/Founder.png',                  TZ1_LEFT, TZ_TOP, TZ_W, TZ_H)
add_image_proportional(s, 'assets/images/about-wicklow.webp', TZ2_LEFT, TZ_TOP, TZ_W, TZ_H)

# ── SLIDE 11 — Community ─────────────────────────────────────────────────────
s = new_slide('Community')
add_label(s, 'SOCIAL PROOF')
add_title(s, "12,000+ people didn't just buy a backpack \u2014 they chose a system.", has_image=True)
add_bullets(s, [
    '6 named archetypes: Aoife (Donegal), Marco (Barcelona), Sarah (Dublin), Tomas (Pamplona), Nadia (Amsterdam), James (London)',
    'These are not testimonials \u2014 they are portraits of people living in the same rhythm',
    '50+ countries. 1,000+ community-tagged photos',
    'Community proves the idea works in real life across different contexts',
    'Basecamp Loyalty Program: the brand invests in long-term relationships, not transactions',
])
add_image_proportional(s, 'CCP/Community.png',   TZ1_LEFT, TZ_TOP, TZ_W, TZ_H)
add_image_proportional(s, 'CCP/Community 1.png', TZ2_LEFT, TZ_TOP, TZ_W, TZ_H)

# ── SLIDE 12 — Sustainability ────────────────────────────────────────────────
s = new_slide('Sustainability')
add_label(s, 'BRAND VALUES')
add_title(s, "Longevity is not a position. It's a construction parameter.", has_image=True)
add_bullets(s, [
    'Product designed for 15 years of use \u2014 not a marketing claim, a design spec',
    'Every wearable component is replaceable: YKK zippers, Duraflex buckles, module connectors',
    'Full materials transparency: denier count, weave type, coating spec, hardware supplier',
    'Modular buyers replace their main pack 73% less often than single-purpose bag buyers',
    '"Where We\'re Still Working" \u2014 an honest section about what hasn\'t been solved yet',
])
add_image_proportional(s, 'CCP/What Goes Into Every Pack sustainability.png', TZ1_LEFT, TZ_TOP, TZ_W, TZ_H)
add_image_proportional(s, 'assets/images/sustainability-repair.webp',         TZ2_LEFT, TZ_TOP, TZ_W, TZ_H)

# ── SLIDE 13 — Visual Language ───────────────────────────────────────────────
s = new_slide('Visual Language')
add_label(s, 'VISUAL SYSTEM')
add_title(s, 'Dark navy. Electric mint. Not sportswear \u2014 cinematic.', has_image=True)
add_bullets(s, [
    'Dark navy #0B1626: not black (mass-market sport), not white (generic shop) \u2014 navy is the TrailVolt world',
    'Electric mint #3FFFAB: one accent source per image, never a background, never a wash',
    'Product photography: like a luxury watch commercial, not an Amazon listing',
    'Lifestyle: person in motion, specific location, editorial photography \u2014 never posed',
    "The Gear Patrol rule: if the image fits their Editor's Picks without a caption \u2014 it's correct",
])
add_image_proportional(s, 'assets/images/pack-everyday-black.webp', TZ1_LEFT, TZ_TOP, TZ_W, TZ_H)
add_image_proportional(s, 'assets/images/about-airport.webp',       TZ2_LEFT, TZ_TOP, TZ_W, TZ_H)

# ── SLIDE 14 — Why Custom ────────────────────────────────────────────────────
s = new_slide('Why Custom')
add_label(s, 'TECHNICAL APPROACH')
add_title(s, "Shopify doesn't know what a modular attachment rail is.", has_image=True)
add_bullets(s, [
    'A boxed platform can do: product card, cart, checkout. Necessary but not sufficient',
    'TrailVolt needed: storytelling \u2192 system explanation \u2192 configurator \u2192 sale',
    'A live configurator with live pricing and cross-compatibility logic cannot exist in a template theme',
    "The brand's visual language does not fit inside any available theme",
    'i18n across 6 languages, PWA with offline mode, schema.org markup \u2014 all written to match the concept',
    'Claude Code + GitHub: not a workaround \u2014 a proper product engineering tool',
])
add_image_proportional(s, 'CCP/languege .png', RZ_LEFT, RZ_TOP, RZ_W, RZ_H)

# ── SLIDE 15 — Final ─────────────────────────────────────────────────────────
s = new_slide('Final')
add_label(s, 'CONCLUSION')
add_title(s, 'TrailVolt is not a site about a backpack.', has_image=True)
add_subtitle(s, 'A digital environment for a product that requires understanding, not just purchasing.', top_cm=3.2)
# large mint statement
txb = s.shapes.add_textbox(Cm(1.8), Cm(4.8), Cm(17), Cm(1.2))
tf = txb.text_frame; p = tf.paragraphs[0]; r = p.add_run()
r.text = 'One Pack. Every Horizon.'; r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = MINT; r.font.name = 'Arial'
add_bullets(s, [
    'The product was born from real frustration \u2014 and every page shows it',
    'The system is explained before it is sold',
    'Premium means depth, honesty, durability and control \u2014 not decoration',
    "The site itself is evidence of the product's standard",
    'Claude Code + GitHub was the right tool for a non-standard problem',
], top_cm=6.5)
add_image_proportional(s, 'assets/images/hero-pack-main.webp', 33.87*0.55, 0, 33.87*0.45, 19.05)

# ── SAVE ─────────────────────────────────────────────────────────────────────
output = 'trailvolt-presentation.pptx'
prs.save(output)

# ── CLEANUP ──────────────────────────────────────────────────────────────────
shutil.rmtree(TMP, ignore_errors=True)

# ── SUMMARY ──────────────────────────────────────────────────────────────────
print('\n\u2500\u2500 Slide summary \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500')
for i, t in enumerate(results, 1):
    print(f'\u2713 Slide {i:02d} \u2014 {t}')
size = os.path.getsize(output)
print(f'\n\u2713 Saved: {output}  ({size:,} bytes)')
